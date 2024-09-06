from transformers import AutoModelForSeq2SeqLM, AutoTokenizer
from utils import split_languages
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed
from threading import Lock

class BasicTranslator():
    def __init__(self, model_name, ignore_terms=None):
        self.model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
        self.tokenizer = AutoTokenizer.from_pretrained(model_name)
        self.ignore_terms = ignore_terms if ignore_terms is not None else []
        self.tgt_lang = ''
        self.src_lang = ''

    def _translate_segment(self, segment):
        translated_segment = ''
        current_pos = 0 

        while current_pos < len(segment):
            next_pos = len(segment)
            for term in self.ignore_terms:
                term_pos = segment.find(term, current_pos)
                if current_pos <= term_pos < next_pos:
                    next_pos = term_pos

            if next_pos == current_pos:
                term = next((term for term in self.ignore_terms if segment.startswith(term, current_pos)), None)
                if term:
                    translated_segment += term
                    current_pos += len(term)
                else:
                    translated_segment += segment[current_pos]
                    current_pos += 1
            else:
                text_to_translate = segment[current_pos:next_pos].rstrip()
                trailing_whitespace = segment[current_pos:next_pos][len(text_to_translate):]
                if text_to_translate:
                    token_ids = self.tokenizer.encode(text_to_translate, return_tensors='pt', padding=True, truncation=True, max_length=128)
                    model_out = self.model.generate(token_ids)
                    output_text = self.tokenizer.decode(model_out[0], skip_special_tokens=True)
                    translated_segment += output_text
                translated_segment += trailing_whitespace
                current_pos = next_pos
                
        return translated_segment
    
    def translate(self, text):
        segments = split_languages(text)
        translated_segments = []
        for segment, lang in segments:
            if lang == self.tgt_lang or lang == 'punct':
                translated_segments.append(segment)
            elif lang == self.src_lang:
                translated_segment = self._translate_segment(segment)
                translated_segments.append(translated_segment)
            elif self.src_lang == '' or self.tgt_lang == '':
                raise ValueError("Source and target languages must be set before translating")
        result_text = ''.join(translated_segments)
        return result_text

    

class OpusEnRuTranslator(BasicTranslator):
    def __init__(self, ignore_terms=None):
        super().__init__("Helsinki-NLP/opus-mt-en-ru", ignore_terms=ignore_terms)
        self.src_lang = 'en'
        self.tgt_lang = 'ru'


class OpusRuEnTranslator(BasicTranslator):
    def __init__(self, ignore_terms=None):
        super().__init__("Helsinki-NLP/opus-mt-ru-en", ignore_terms=ignore_terms)
        self.src_lang = 'ru'
        self.tgt_lang = 'en'


class PptxTranslator():
    def __init__(self, src_lang, tgt_lang, ignore_terms=None):
        if src_lang == "en" and tgt_lang == "ru":
            self.translator = OpusEnRuTranslator(ignore_terms=ignore_terms)
        elif src_lang == "ru" and tgt_lang == "en":
            self.translator = OpusRuEnTranslator(ignore_terms=ignore_terms)

    def translate_shape_text(self, shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for subshape in shape.shapes:
                self.translate_shape_text(subshape)
        elif hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():  # Only translate non-empty strings
                        translated_text = self.translator.translate(text=run.text)
                        run.text = translated_text

    def translate_slide(self, slide):
        shapes = slide.shapes
        with ThreadPoolExecutor(max_workers=8) as shape_executor:
            futures = [shape_executor.submit(self.translate_shape_text, shape) for shape in shapes]
    
    def translate_presentation(self, original_file_path, translated_file_path, progress_callback=None, finished_callback=None):
        # Make a copy of the entire presentation
        shutil.copy2(original_file_path, translated_file_path)
        
        prs = Presentation(translated_file_path)
        total_slides = len(prs.slides)
        slides_processed_lock = Lock()
        slides_processed = 0
        
        with ThreadPoolExecutor(max_workers=8) as slide_executor:
            futures = [slide_executor.submit(self.translate_slide, slide) for slide in prs.slides]
            
            for future in as_completed(futures):
                try:
                    with slides_processed_lock:
                        slides_processed += 1
                    progress = int((slides_processed / total_slides) * 100)
                    if progress_callback:
                        progress_callback(progress)
                except Exception as e:
                    print(f"An error occurred while translating a slide: {e}")
            future.result()
            if finished_callback:
                finished_callback()

        prs.save(translated_file_path)
